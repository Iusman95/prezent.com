from django.shortcuts import render, redirect, get_object_or_404
from prezents.forms import DowlandFileForm, CreatePostForm, CommentForm
from user_profile.models import Files
from prezents.models import Post, Comment
from django.urls import reverse




#Представление для вывода постов
def index(request):
    posts = Post.objects.all()
    return render(request, 
        template_name='prezents/index.html',
        context={'posts': posts})



def detail(request, pk):
    post =  get_object_or_404(Post, pk=pk)
    comment_form = CommentForm()
    comments = Comment.objects.filter(post=post)

    context = {
        'form': comment_form,
        'post': post,
        'comments': comments
    }

    return render(request, 
            template_name='prezents/detail.html', 
            context=context)





#Представление для страницы "my_files"
def my_files(request):
    files = Files.objects.filter(user=request.user)
    return render(request, 
        template_name='prezents/my_files.html', 
        context={'files': files})




#Представление для загрузки фотографий на профиль
def add_image_on_profile(request):
    if request.method == "POST":
        form_two = DowlandFileForm(request.POST, request.FILES)
        if form_two.is_valid():
            form_two.instance.user = request.user
            form_two.save()
            return redirect('my_files')
    else:
        form_two = DowlandFileForm()
    
    return render(request, 
        template_name='prezents/upload_file.html', 
        context={'form_two': form_two})




#Представление для сохранение поста
def add_post(request):
    if request.method == "POST":
        form_ = CreatePostForm(request.POST, request.FILES)
        if form_.is_valid():
            form_.instance.user = request.user
            form_.save()
            return redirect('index')
    else:
        form_ = CreatePostForm()
    
    return render(request, 
        template_name='prezents/upload_index.html', 
        context={'form_': form_})


def add_comment(request, pk):

    if request.method == "POST":
        form = CommentForm(request.POST)
        post = get_object_or_404(Post, id=pk)
        if form.is_valid():
            form.save(commit=False)
            form.instance.user = request.user
            form.instance.post = post
            form.save()
            return redirect(reverse('detail', kwargs={'pk': pk}))
            



















""" from pptx import Presentation  

def prezent(request):
    root = Presentation() #создание объекта презентации
    first_slide_layout = root.slide_layouts[0] # создание макета слайда
    slide = root.slides.add_slide(first_slide_layout) # Создание объекта слайда для добавления  в ppt, т.е. прикрепление слайдов с презентацией, т.е. ppt
    slide.shapes.title.text = "Привет"
    slide.placeholders[1].text = " This is 2nd way"
    root.save("Output.pptx")
    print("done")
    return render(request, 
            template_name='prezents/prezent.html', 
            context={'root': root})  """